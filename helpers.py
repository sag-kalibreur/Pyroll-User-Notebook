from io import StringIO, BytesIO
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.animation as animation

from pptx.util import Inches, Pt

import pyroll.core as pr
from pyroll.from_dict import from_dict as pr_from_dict

from UnitConversion import to_convert, convert_SI

pyroll_possible_imports = [
                      'pyroll.core',
                      'pyroll.freiberg_flow_stress',
                      'pyroll.wusatowski_spreading',
                      'pyroll.marini_spreading',
                      'pyroll.hill_spreading',
                      'pyroll.lendl_equivalent_method',
                      'pyroll.lippmann_mahrenholz_force_torque',
                      'pyroll.zouhar_contact',
                      'pyroll.integral_thermal',
                      'pyroll.ring_model_thermal',
                     ]

plot_settings = {

    'in_label': 'in_profile',
    'out_label': 'out_profile',

    'in_color': 'red',
    'out_color': 'blue',
    'tool_color': 'grey',
    'alpha': 0.5,

    # bar-plot
    'kind': 'barh',
    'rot': 0,
    'width': 0.9,
    'inch_cm': 2.54
    }

cm = 1/2.54

# 16:9

pic = {
#      'vpos_empty' : 0.0,
#      'hpos_empty' : 0.0,
       'vpos_empty' : 3.6,
       'hpos_empty' : 1.0,
       'vpos': 5.86,
       'hpos1': 1.33,
       'hpos2': 18.5,
       'height': 9.0,
       'width': 17.0,
#      'height_max': 19.05,
#      'width_max': 33.87,
       'height_max': 13.5,
       'width_max': 32.0
}

pic = {key : value*cm for key, value in pic.items()}

tab = {'vpos1': 4.41,
       'vpos2': 14.15,
       'hpos': pic['hpos1'],
       'height1': 9.35,
       'height2': 2.,
       'width': 32.0}

tab = {key : value*cm for key, value in tab.items()}
tab['hpos'] = pic['hpos1']

def GetProfileData(DB, profile_id):
    
    df_profiles = pd.read_excel(DB['xlsx_file'], sheet_name=DB['sheet'])

    geometry_types = [
                    'diameter',
                    'height',
                    'width', 
                    'side',
                    'corner_radius'
                    ]

    ProfileDB = df_profiles.set_index('ID')

    profile_type = (ProfileDB
                    .query('ID == @profile_id')
                    .filter(['profile_type'])
                    .to_dict('index')[profile_id]
                    )

    geometry_data = (ProfileDB
                    .query('ID == @profile_id')
                    .dropna(axis=1, how='all')
                    .filter(geometry_types)
                    .to_dict('index')[profile_id]
                    )

    return profile_type, geometry_data

def CreateProfile(parameters):

    profile_type = parameters.pop('profile_type')

    match profile_type:
        case 'square':
            profile = pr.Profile.square(**parameters)
        case 'round':
            profile = pr.Profile.round(**parameters)
        case 'box':
            profile = pr.Profile.box(**parameters)
        case 'hex':
            profile = pr.Profile.hexagon(**parameters)
        case _:
            print('profile_type {} not implemented!').format(profile_type)

    return profile

def SetupStandDefinition(Sequence, rollDB, COULOMB_FRICTION, dimension_ctrl, solve_data):
    
    for roll_name in Sequence:
    
        PassData = rollDB[roll_name]
    
        rollpass_type = PassData.pop('type')
        nominal_diameter = PassData.pop('nominal_diameter')
        
        roll = {
            '__ctor__': 'Roll',
            'groove': pr.create_groove_by_type_name(**PassData),
            'nominal_diameter': nominal_diameter,
            }
        
        roller = {
            '__ctor__': rollpass_type,
            'label': roll_name,
            'roll': roll,
            }
    
        try:
            PassSequence.append(pr_from_dict(roller, {}))
        except:
            PassSequence = [pr_from_dict(roller, {})]
    
    for n, Stand in enumerate(PassSequence):

        if solve_data['solve_type'] in ['given_RPM']:
            RPM = solve_data['RPM']
            Stand.roll.rotational_frequency = RPM[n]

        if solve_data['solve_type'] in ['given_vel']:
            vel = solve_data['v_target']
            Stand.velocity = vel[n]
            
        Stand.coulomb_friction_coefficient = COULOMB_FRICTION[n]
    
        if Stand.__str__().startswith('TwoRollPass'):
            try:
                StandCounter += 1
            except:
                StandCounter = 1
    
            Stand.gap = dimension_ctrl[n]
            Stand.orientation = 'v' if StandCounter%2 == 0 else 'h'
    
        elif Stand.__str__().startswith('ThreeRollPass'):
            try:
                KocksCounter += 1
            except:
                KocksCounter = 1
    
            Stand.inscribed_circle_diameter = dimension_ctrl[n]
            Stand.orientation = 'AntiY' if KocksCounter%2 == 0 else 'Y'

    return PassSequence
    
def SetupStandDefinition_V(Sequence, rollDB, COULOMB_FRICTION, dimension_ctrl, solve_data):
    
    for roll_name in Sequence:
    
        PassData = rollDB[roll_name]
    
        rollpass_type = PassData.pop('type')
        nominal_diameter = PassData.pop('nominal_diameter')
        
        roll = {
            '__ctor__': 'Roll',
            'groove': pr.create_groove_by_type_name(**PassData),
            'nominal_diameter': nominal_diameter,
            }
        
        roller = {
            '__ctor__': rollpass_type,
            'label': roll_name,
            'roll': roll,
            }
    
        try:
            PassSequence.append(pr_from_dict(roller, {}))
        except:
            PassSequence = [pr_from_dict(roller, {})]
    
    for n, Stand in enumerate(PassSequence):

        if solve_data['solve_type'] in ['given_RPM']:
            RPM = solve_data['RPM']
            Stand.roll.rotational_frequency = RPM[n]

        if solve_data['solve_type'] in ['given_vel']:
            vel = solve_data['v_target']
            Stand.velocity = vel[n]
            
        Stand.coulomb_friction_coefficient = COULOMB_FRICTION[n]
    
        if Stand.__str__().startswith('TwoRollPass'):
            try:
                StandCounter += 1
            except:
                StandCounter = 1
    
            Stand.gap = dimension_ctrl[n]
            Stand.orientation = 'h' if StandCounter%2 == 0 else 'v'
    
        elif Stand.__str__().startswith('ThreeRollPass'):
            try:
                KocksCounter += 1
            except:
                KocksCounter = 1
    
            Stand.inscribed_circle_diameter = dimension_ctrl[n]
            Stand.orientation = 'AntiY' if KocksCounter%2 == 0 else 'Y'

    return PassSequence

def AddTransports(PassSequence, transport_lengths):

    Transports = [pr.Transport(label='Transport_{}'.format(n), length=length) for n, length in enumerate(transport_lengths, 1)]
    
    PassIndexes = [n for n, _ in enumerate(np.empty_like(PassSequence))]
    TranspIndexes = [i+len(PassIndexes) for i in PassIndexes]
    
    NewIndexes = list(sum(zip(PassIndexes, TranspIndexes), ()))[:-1]
    
    PassSequence += Transports
    
    return [PassSequence[i] for i in NewIndexes]
    
def AddTransports(PassSequence, durations):

    Transports = [pr.Transport(label='Transport_{}'.format(n), length=length) for n, length in enumerate(durations, 1)]
    
    PassIndexes = [n for n, _ in enumerate(np.empty_like(PassSequence))]
    TranspIndexes = [i+len(PassIndexes) for i in PassIndexes]
    
    NewIndexes = list(sum(zip(PassIndexes, TranspIndexes), ()))[:-1]
    
    PassSequence += Transports
    
    return [PassSequence[i] for i in NewIndexes]

def SolveSequence(sequence, in_profile, solve_data):
    
    match solve_data['solve_type']:
        case 'forward':
            print(solve_data['solve_type'])
            sequence.solve_velocities_forward(in_profile=in_profile,
                                             initial_speed=solve_data['v_target'])
        case 'backward':
            print(solve_data['solve_type'])
            sequence.solve_velocities_backward(in_profile=in_profile,
                                        final_speed=solve_data['v_target'],
                                        final_cross_section_area=sequence[-1].usable_cross_section.area)
        case 'given_RPM' | 'given_vel':
            print(solve_data['solve_type'])
            sequence.solve(in_profile)

        case _:
            print('Strategy {} not implemented'.format(solve_data['solve_type']))

def CreateRollDB(ToolDB, convert_SI):
    
    df_tools = pd.read_excel(ToolDB['xlsx_file'], sheet_name=ToolDB['sheet'])
    
    df_tools[[
        'r1', 'r2', 'r3',
        'depth', 'usable_width', 'ground_width',
        'nominal_diameter', 
        'pad'
       ]] *= convert_SI['length']

    RollDB = (df_tools
              .set_index('label')
              .to_dict('index')
             )
    
    # remove NaN-entries
    for key, dic in RollDB.items():
        RollDB[key] = {k: dic[k] for k in dic if type(dic[k]) is str or not np.isnan(dic[k])}

    return RollDB

def Modify_df(df, ):

    df['Source'] = 'PyRolL'
    df['StandPos'] = df.index + 1
    
    df['roll_entry_angle'] *= -1 # negative entry_angle is unintuitive...
              
## alles nur um Konturen für die Anzeige aufzubereiten...

    results = [
        'in_profile_cross_section',
        'out_profile_cross_section',
        'in_profile_lendl_section',
        'out_profile_lendl_section',
        'usable_cross_section',
        'displaced_cross_section',
        'reappearing_cross_section',
        'in_profile_equivalent_rectangle',
        'out_profile_equivalent_rectangle',
        ]
    
    for result in results:

        df.loc[df['orientation']=='AntiY', result+'_y'] *= -1
        
        swap_x = df.loc[df['orientation']=='v', result+'_x'].copy().values
        swap_y = df.loc[df['orientation']=='v', result+'_y'].copy().values

#        print(result, '\n', df.loc[ df['orientation']=='v', [result+'_x']], '\n', swap_y)
        
        df.loc[df['orientation']=='v', result+'_x'] = -1 * swap_y
        df.loc[df['orientation']=='v', result+'_y'] = swap_x

## sollte bis UM überflüssing sein
    
    for unit, entry_list in to_convert.items():
        for entry in entry_list:

            if entry in df.columns:
                if unit == 'temperature':
                    df[entry] -= convert_SI[unit]
                elif unit == 'angle':
                    df[entry] = np.degrees(df[entry])                   
                else:
                    df[entry] /= convert_SI[unit]

    return df

def AnimateSequence(df, limit, plot_settings, size):
    
    fig = plt.figure(figsize=(size, size))
    
    ax = plt.axes(xlim=(-limit, limit), ylim=(-limit, limit))
    ax.set_aspect('equal')
    
    section1, = ax.plot([], [], ':', c=plot_settings['tool_color'], lw=1)
    section2, = ax.plot([], [], lw=1, c=plot_settings['in_color'], label=plot_settings['in_label'])
    section3, = ax.plot([], [], lw=1, c=plot_settings['out_color'], label=plot_settings['out_label'])
    
    dfPyRolL = df.query("Source == 'PyRolL'")
    
    def update(frame):
        
        ax.set_title(u"Standposition = {}".format(dfPyRolL.label[frame]))
        
        section1.set_xdata(dfPyRolL.technologically_orientated_contour_lines_x[frame])
        section1.set_ydata(dfPyRolL.technologically_orientated_contour_lines_y[frame])
    
        section2.set_xdata(dfPyRolL.in_profile_cross_section_x[frame])
        section2.set_ydata(dfPyRolL.in_profile_cross_section_y[frame])
    
        section3.set_xdata(dfPyRolL.out_profile_cross_section_x[frame])
        section3.set_ydata(dfPyRolL.out_profile_cross_section_y[frame])
        
        return()
    
    ani = animation.FuncAnimation(fig=fig, func=update, frames=len(df['StandPos'].unique()), interval=500)
    plt.legend()
    plt.close()

    return ani

def axplot(dfplot, idx, prop, kind, ax, rot, width, annotate):

    if kind in ['bar', 'barh']:
        plot_opts = {'kind': kind, 'ax': ax, 'rot': rot, 'width': width}

    elif kind == 'line':
        plot_opts = {'kind': kind, 'ax': ax, 'rot': rot}


    dfplot.set_index(idx)[prop].plot(**plot_opts)

    if isinstance(prop, list):
        prop = prop[0]
        title = annotate[prop]['title']
        xlabel = annotate[prop]['unit']
        multi = True
    else:
        title = prop
        xlabel = annotate[prop]['unit']
        multi = False

    for container in ax.containers:
        ax.bar_label(container,
                     fmt=annotate[prop]['fmt'],
                     label_type=annotate[prop]['label'],
                     color=annotate[prop]['font_color'],
                     fontsize=7 if multi else 9)

    ax.set_title(title)

    if kind in ['bar', 'barh']:

        ax.set_xlabel(xlabel)
        ax.set_ylabel(None)
        ax.invert_yaxis()

    elif kind == 'line':

        ax.set_xlabel(None)
        ax.set_ylabel(None)

def AddPictureToSlide(prs, layouts, slide_type, pic, fig, slide_title=None):

    slide = prs.slides.add_slide(layouts[slide_type])
    image_stream = BytesIO()
    fig.savefig(image_stream)

    if slide_type == 'Title_Only':
        slide.shapes.title.text = slide_title
        
    slide.shapes.add_picture(image_stream, Inches(pic['hpos_empty']), Inches(pic['vpos_empty']))

def ComparisonPlot(df, comp_results, limit, prs, layouts, slide_type, slide_title=None, filled=False, eq_rec=False):
    
    for pos in df['StandPos'].unique():
        
        fig, ((ax1), (ax4)) = plt.subplots(
                                            nrows=1, ncols=2,
#                                            layout="constrained"
                                            )

        slide_label = df.query("Source == 'PyRolL' & StandPos == @pos")['label'].values[0]
    
        x_in = df.query("Source == 'PyRolL' & StandPos == @pos")[comp_results[0]+'_x'].explode()
        y_in = df.query("Source == 'PyRolL' & StandPos == @pos")[comp_results[0]+'_y'].explode()
        
        x_out = df.query("Source == 'PyRolL' & StandPos == @pos")[comp_results[1]+'_x'].explode()
        y_out = df.query("Source == 'PyRolL' & StandPos == @pos")[comp_results[1]+'_y'].explode()
    
        x_roll = df.query("Source == 'PyRolL' & StandPos == @pos").technologically_orientated_contour_lines_x.explode()
        y_roll = df.query("Source == 'PyRolL' & StandPos == @pos").technologically_orientated_contour_lines_y.explode()

        x_in_eq = df.query("Source == 'PyRolL' & StandPos == @pos").in_profile_equivalent_rectangle_x.explode()
        y_in_eq = df.query("Source == 'PyRolL' & StandPos == @pos").in_profile_equivalent_rectangle_y.explode()
 
        x_out_eq = df.query("Source == 'PyRolL' & StandPos == @pos").out_profile_equivalent_rectangle_x.explode()
        y_out_eq = df.query("Source == 'PyRolL' & StandPos == @pos").out_profile_equivalent_rectangle_y.explode()

        if comp_results[0].startswith('usable'):
            color=plot_settings['tool_color']
            label='usable'
        else:
            color=plot_settings['in_color']
            label=plot_settings['in_label']
        
        if filled:
            ax1.fill(x_in, y_in, color=color, label=label)
            ax1.fill(x_out, y_out, color=plot_settings['out_color'], alpha=plot_settings['alpha'], label=plot_settings['out_label'])
            ax1.plot(x_roll, y_roll, ':', color=plot_settings['tool_color'])
            ax4.fill(x_in, y_in, color=color, label=plot_settings['in_label'])
            ax4.fill(x_out, y_out, color=plot_settings['out_color'], alpha=plot_settings['alpha'], label=plot_settings['out_label'])

        else:
            ax1.plot(x_in, y_in, color=plot_settings['in_color'], label=plot_settings['in_label'])
            ax1.plot(x_out, y_out, color=plot_settings['out_color'], label=plot_settings['out_label'])
            ax1.plot(x_roll, y_roll, ':', color=plot_settings['tool_color'])
            ax4.plot(x_in, y_in, color=plot_settings['in_color'], label=plot_settings['in_label'])
            ax4.plot(x_out, y_out, color=plot_settings['out_color'], label=plot_settings['out_label'])
        
        if eq_rec:
            ax1.plot(x_in_eq, y_in_eq, ':', color=plot_settings['in_color'], label=plot_settings['in_label'])
            ax1.plot(x_out_eq, y_out_eq, ':', color=plot_settings['out_color'], label=plot_settings['out_label'])
            ax4.plot(x_in_eq, y_in_eq, ':', color=plot_settings['in_color'], label=plot_settings['in_label'])
            ax4.plot(x_out_eq, y_out_eq, ':', color=plot_settings['out_color'], label=plot_settings['out_label'])

        ax1.set_xlim(-limit, limit)
        ax1.set_ylim(-limit, limit)
        ax1.text(0.05, 0.95, slide_label, transform=ax1.transAxes)
        ax1.legend(title=comp_results[2])
        
        ax1.set_aspect('equal')
        ax4.set_aspect('equal')

        AddPictureToSlide(prs, layouts, slide_type, pic, fig, '{} {}'.format(slide_title, slide_label))

def ComparisonPlot_polar(df, comp_results, rmax, prs, layouts, slide_type, slide_title=None):
    
    for pos in df['StandPos'].unique():
    
        fig = plt.figure(
#                        layout="constrained"
                        )
    
        ax1 = fig.add_subplot(121, projection='polar')
    
        ax4 = fig.add_subplot(122)

        slide_label = df.query("Source == 'PyRolL' & StandPos == @pos")['label'].values[0]
    
        radius, theta = df.query("Source == 'PyRolL' & StandPos == @pos")[comp_results].explode()
        
        ax1.plot(theta, radius, c=plot_settings['out_color'])
        ax1.set_rmax(rmax)
    
    
        ax4.plot(theta, radius, c=plot_settings['out_color'])
        ax4.set_xticks([-np.pi, -np.pi/2, 0, np.pi/2, np.pi], labels=['180', '270', '0', '90', '180'])

        AddPictureToSlide(prs, layouts, slide_type, pic, fig, '{} {}'.format(slide_title, slide_label))

def ContactAreaPlot(df, comp_results, limit, prs, layouts, slide_type, slide_title=None):
    
    for pos in df['StandPos'].unique():
    
        fig = plt.figure(
    #                    layout="constrained"
                        )
    
        ax1 = fig.add_subplot(121, projection='3d')
    
        ax4 = fig.add_subplot(122)

        slide_label = df.query("Source == 'PyRolL' & StandPos == @pos")['label'].values[0]
        
        dftmp = df.loc[(df['Source'] == 'PyRolL') & (df['StandPos'] == pos)]
    
        zouhar_coords = np.array([
            [-dftmp.out_profile_width[pos - 1] / 2., 0],
            [-dftmp.out_profile_width[pos - 1] / 2, dftmp.roll_contact_length[pos - 1] * dftmp.zouhar_contact_c2[pos - 1]],
            [-dftmp.zouhar_contact_in_width[pos - 1] *  dftmp.zouhar_contact_c1[pos - 1] / 2, dftmp.roll_contact_length[pos - 1]],
            [ dftmp.zouhar_contact_in_width[pos - 1] *  dftmp.zouhar_contact_c1[pos - 1] / 2, dftmp.roll_contact_length[pos - 1]],
            [ dftmp.out_profile_width[pos - 1] / 2, dftmp.roll_contact_length[pos - 1] * dftmp.zouhar_contact_c2[pos - 1]],
            [ dftmp.out_profile_width[pos - 1] / 2, 0],
            [-dftmp.out_profile_width[pos - 1] / 2, 0]
                        ])
    
        x_in = df.query("Source == 'PyRolL' & StandPos == @pos").in_profile_cross_section_x.explode()
        y_in = df.query("Source == 'PyRolL' & StandPos == @pos").in_profile_cross_section_y.explode()
        
        x_out = df.query("Source == 'PyRolL' & StandPos == @pos").out_profile_cross_section_x.explode()
        y_out = df.query("Source == 'PyRolL' & StandPos == @pos").out_profile_cross_section_y.explode()
    
        ax1.set_xlim(-110, 110)
        ax1.set_ylim(0, 120)
        ax1.set_zlim(-110, 110)
        ax1.set_aspect('equal')
    
        ax1.plot(zouhar_coords[:,0], zouhar_coords[:,1])
        
        ax1.plot(
                x_in,
                np.zeros_like(x_in),
                y_in,
                color=plot_settings['in_color'], label=plot_settings['in_label'])
    
        ax1.plot(
                x_out,
                np.zeros_like(x_out),
                y_out,
                color=plot_settings['out_color'], label=plot_settings['out_label'])
        
        ax4.set_xlim(-limit, limit)
        ax4.set_ylim(-limit, limit)
        ax4.set_aspect('equal')
        
        ax4.plot(x_in,
                 y_in,
                 ':', color=plot_settings['in_color'], label=plot_settings['in_label'])
        ax4.plot(x_out,
                 y_out,
                 ':', color=plot_settings['out_color'], label=plot_settings['out_label'])
    
        ax4.fill(zouhar_coords[:,0], zouhar_coords[:,1], alpha=0.5)
        
        ax4.text(0.05, 0.95, slide_label, transform=ax4.transAxes)
        ax4.legend()

        AddPictureToSlide(prs, layouts, slide_type, pic, fig, '{} {}'.format(slide_title, slide_label))

def TablePlot(df, prs, layouts, slide_type, slide_title=None):
    
    df_table = df[[
                    'StandPos',
                    'StandPos_line',
                    'label',
                    'roll_nominal_diameter',
                    'roll_working_diameter',
                    'coulomb_friction_coefficient',
                    'GT',
                    'DT',
                    'dim_ctrl_type',
                    'dim_ctrl',
                    'out_profile_cross_section_area',
                    'reduction',
                    'elongation',
                    'spread',
                    'roll_entry_angle',
                    'gear_levels',
                    'motor_rpms',
                    'velocity'
                    ]].copy().astype(str)
    
    df_table.update(df_table[['motor_rpms', 'out_profile_cross_section_area']].map(lambda x: f' {float(x):.0f} '))
    
    df_table.update(df_table[[
                            'roll_nominal_diameter',
                            'roll_working_diameter',
     #                       'DT',
     #                       'GT',
                            'dim_ctrl',
                            'reduction',
                            'roll_entry_angle',
                            ]].map(lambda x: f' {float(x):.1f} '))
    
    df_table.update(df_table[['coulomb_friction_coefficient', 'elongation', 'spread', 'DT', 'GT', 'velocity']].map(lambda x: f' {float(x):.2f} '))
    
    df_table.update(df_table[['dim_ctrl_type', 'label']].map('  {:s}  '.format))
    
    fig, ax = plt.subplots(
                           layout="constrained"
                          )
    
    better_labels = {
                    'StandPos': 'Pass',
                    'StandPos_line': 'Stand',
                    'label': 'Label',
                    'roll_nominal_diameter': 'Roll-Ø',
                    'roll_working_diameter': 'Working-Ø',
                    'coulomb_friction_coefficient' : 'µ',
                    'GT': 'H',
                    'DT': 'W',
                    'roll_entry_angle': 'Entry angle',
                    'dim_ctrl': '',
                    'dim_ctrl_type': '',
                    'out_profile_cross_section_area': 'Area',
                    'reduction': 'Reduction',
                    'elongation': 'Elongation',
                    'spread': 'Spread',
                    'gear_levels': 'Gear',
                    'motor_rpms': 'Motor RPM',
                    'velocity': 'Velocity',
                    }
    
    ax.axis('off')
    
    for key in df_table.keys():
        if key in better_labels:
            df_table = df_table.rename(columns={key: better_labels[key]})
    
    tab2plot = ax.table(cellText=df_table.values,
                        colLabels=df_table.keys(),
                        loc='center').auto_set_column_width(col=list(range(len(df_table.columns))))

    AddPictureToSlide(prs, layouts, slide_type, pic, fig, slide_title)

def AddChartSlide(Sequence, df, ChartSlides, prs, layouts):

    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.chart.data import ChartData
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

    for key, result in ChartSlides.items():

        data = df[result]

        rows = 2
        cols = len(data)
        
        slide = prs.slides.add_slide(layouts['Title_Only'])

        slide.shapes.title.text = '{} Verlauf'.format(key)

        shape = slide.shapes.add_table(rows, cols,
                                   Inches(tab['hpos']), Inches(tab['vpos2']),
                                   Inches(tab['width']), Inches(tab['height2']))
    
        for i, datavalue in enumerate(data):
            shape.table.cell(0, i).text = '{}'.format(Sequence[i])
            shape.table.cell(1, i).text = '{:.1f}'.format(datavalue)

        def iter_cells(table):
            for row in table.rows:
                for cell in row.cells:
                    yield cell

        for cell in iter_cells(shape.table):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(11)

        chart_data = ChartData()
                
        chart_data.categories = Sequence
    
        chart_data.add_series(key, np.array(data))
        
        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                               Inches(tab['hpos']),
                               Inches(tab['vpos1']),
                               Inches(tab['width']),
                               Inches(tab['height1']),
                               chart_data).chart
        
        chart.has_legend = False
        chart.has_title = True
        chart.plots[0].vary_by_categories = False

        from pptx.oxml.xmlchemy import OxmlElement

        shape_properties = OxmlElement("c:spPr")
        chart.element.append(shape_properties)
        fill_properties = OxmlElement("a:solidFill")
        shape_properties.append(fill_properties)
        scheme_color = OxmlElement("a:schemeClr")
        color_value = dict(val="bg2")
        scheme_color.attrib.update(color_value)
        fill_properties.append(scheme_color)

        #rgb_color = OxmlElement("a:srgbClr")
        #color_value = dict(val='%02x%02x%02x' % (130, 130, 130))
        #rgb_color.attrib.update(color_value)
        #fill_properties.append(rgb_color)
        
        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS_STACKED,
                               Inches(tab['hpos']),
                               Inches(tab['vpos1']),
                               Inches(tab['width']),
                               Inches(tab['height1']),
                               chart_data).chart
        
        chart.has_legend = False
        chart.has_title = True
